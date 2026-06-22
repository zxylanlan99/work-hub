"""
文件解析模块 — 支持 PDF / PPT / Word / Markdown / TXT
提取纯文本内容，供后续切片和向量化使用
"""
import os
import logging

logger = logging.getLogger(__name__)


def parse_file(file_path: str, file_name: str = None) -> str:
    """
    根据文件扩展名自动选择解析器，提取纯文本

    Args:
        file_path: 文件在服务器上的绝对路径
        file_name: 原始文件名（用于判断扩展名，默认从 file_path 提取）

    Returns:
        提取出的纯文本字符串

    Raises:
        ValueError: 不支持的文件格式
        Exception: 解析过程中的错误
    """
    name = file_name or os.path.basename(file_path)
    ext = name.rsplit('.', 1)[-1].lower() if '.' in name else ''

    logger.info(f"开始解析文件: {name} (格式: {ext})")

    if ext in ('md', 'markdown', 'txt'):
        text = _parse_text(file_path)
    elif ext == 'pdf':
        text = _parse_pdf(file_path)
    elif ext == 'docx':
        text = _parse_docx(file_path)
    elif ext == 'doc':
        raise ValueError('.doc 旧格式不支持，请转换为 .docx 后上传')
    elif ext == 'pptx':
        text = _parse_pptx(file_path)
    elif ext == 'ppt':
        raise ValueError('.ppt 旧格式不支持，请转换为 .pptx 后上传')
    else:
        # 未知格式，尝试当文本读
        logger.warning(f"未知文件格式: {ext}，尝试作为文本读取")
        text = _parse_text(file_path)

    if not text or len(text.strip()) < 10:
        raise ValueError('文件内容为空或无法提取文本')

    logger.info(f"文件解析完成: {name}, 提取文本 {len(text)} 字符")
    return text.strip()


def _parse_text(file_path: str) -> str:
    """读取纯文本文件 (MD / TXT)"""
    encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    # 最后兜底
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()


def _parse_pdf(file_path: str) -> str:
    """使用 PyMuPDF (fitz) 解析 PDF"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    full_text = []
    for page_num in range(doc.page_count):
        page = doc[page_num]
        page_text = page.get_text("text")
        if page_text.strip():
            full_text.append(page_text.strip())
    doc.close()

    return '\n\n---\n\n'.join(full_text)


def _parse_docx(file_path: str) -> str:
    """使用 python-docx 解析 Word .docx"""
    from docx import Document

    doc = Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # 提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells)
            if row_text.strip(' |'):
                paragraphs.append(row_text)

    return '\n\n'.join(paragraphs)


def _parse_pptx(file_path: str) -> str:
    """使用 python-pptx 解析 PowerPoint .pptx"""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(' |'):
                        slide_texts.append(row_text)
        if slide_texts:
            slides_text.append(f'## 幻灯片 {slide_num}\n' + '\n'.join(slide_texts))

    return '\n\n---\n\n'.join(slides_text)
